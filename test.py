from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MSO_SHAPE

'''
1. Remove placeholders in slide
https://python-pptx.readthedocs.io/en/latest/user/placeholders-understanding.html
https://stackoverflow.com/questions/39603318/how-to-delete-unpopulated-placeholder-items-using-python-pptx


'''

prs = Presentation()


def slide0():
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Hello, World!"
    subtitle.text = "python-pptx was here!"
    prs.save('test.pptx')


def slide1():
    bullet_slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Adding a Bullet Slide'

    tf = body_shape.text_frame
    tf.text = 'Find the bullet slide layout'

    p = tf.add_paragraph()
    p.text = 'Use _TextFrame.text for first bullet'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
    p.level = 2

    prs.save('test.pptx')


def slide2():
    img_path = 'res/kongfu.jpg'

    blank_slide_layout = prs.slide_layouts[2]
    slide = prs.slides.add_slide(blank_slide_layout)

    # prs.slide_layouts.remove(blank_slide_layout)
    # slide = prs.slides.add_slide(blank_slide_layout)

    # print(slide.shapes._spTree)
    # for i, j in enumerate(slide.shapes._spTree):
    #     print(i, ', ', j)
    # _shape0 = slide.shapes._spTree[0]
    # _shape1 = slide.shapes._spTree[1]
    # slide.shapes._spTree.remove(_shape0)
    # slide.shapes._spTree.remove(_shape1)

    clean_default_placeholders(slide)

    left = top = Inches(1)
    height = Inches(1)
    pic = slide.shapes.add_picture(img_path, left, top, height=height)

    left = Inches(5)
    height = Inches(1.5)
    pic = slide.shapes.add_picture(img_path, left, top, height=height)

    prs.save('test.pptx')


def draw_line(shape):
    line = shape.line
    line.color.rgb = RGBColor(0, 0, 0)
    line.color.brightness = 1
    line.width = Pt(2.5)


def clean_default_placeholders(slide):
    textbox = slide.shapes[0]
    sp = textbox.element
    sp.getparent().remove(sp)
    textbox = slide.shapes[0]
    sp = textbox.element
    sp.getparent().remove(sp)


if __name__ == "__main__":
    # execute only if run as a script
    slide0()
    slide1()
    slide2()
