import sys, getopt
import time
import random
import math
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT as PP_ALIGN
from pptx.enum.text import MSO_VERTICAL_ANCHOR as MSO_ANCHOR
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MSO_SHAPE
from pptx.enum.shapes import MSO_CONNECTOR_TYPE as MSO_CONNECTOR
from pptx.enum.dml import MSO_FILL_TYPE as MSO_FILL
from pptx.enum.dml import MSO_LINE_DASH_STYLE as MSO_LINE

'''
1. Remove placeholders in slide
https://python-pptx.readthedocs.io/en/latest/user/placeholders-understanding.html
https://stackoverflow.com/questions/39603318/how-to-delete-unpopulated-placeholder-items-using-python-pptx

2. Specifies the horizontal alignment for one or more paragraphs.
https://python-pptx.readthedocs.io/en/latest/api/text.html#paragraph-objects
https://python-pptx.readthedocs.io/en/latest/api/enum/PpParagraphAlignment.html#ppparagraphalignment

'''

"""
Side Rect (Inches):
       10
+--------------+
|              |    7.5 
|              |
+--------------+
Width & Height of one slide is about 10 & 7.5, and width : height is 4 : 3

"""

# --------------------------------------------------------------
# custom variables
# --------------------------------------------------------------
# for TEXT type ROWS,COLS => 4, 5
# for GOODS type ROWS, COLS => 3, 4
ROWS = 4
COLS = 5
RESULT_LOWER_LIMIT = 1
PAGE_CONTENT_STYLE = [
    {'op': 'minus', 'result_on_top': True, 'result_upper_limit': 10},
    {'op': 'add', 'result_on_top': False, 'result_upper_limit': 10},
    {'op': 'minus', 'result_on_top': False, 'result_upper_limit': 10},
    {'op': 'add', 'result_on_top': True, 'result_upper_limit': 10},
    {'op': 'minus', 'result_on_top': True, 'result_upper_limit': 10},
    {'op': 'add', 'result_on_top': False, 'result_upper_limit': 10},
    #{'op': 'minus', 'result_on_top': True, 'result_upper_limit': 10},
    #{'op': 'add', 'result_on_top': False, 'result_upper_limit': 10},
    #{'op': 'minus', 'result_on_top': True, 'result_upper_limit': 10},
    #{'op': 'add', 'result_on_top': False, 'result_upper_limit': 10},
]

# Weight of random number
# Rules:
# - First element means the min number.
# - Last element means the max number
# - Element weight uses the the previous weight of element
RESULT_WEIGHT = {
    0: 1,
    #1: 1,
    #2: 1,
    3: 2,
    4: 3,
    #5: 3,
    6: 4,
    7: 6,
    #8: 5,
    #9: 5,
    10: 6,
    #20: 8
}

# --------------------------------------------------------------
# CONST VARIABLES. DO NOT MODIFY!
# --------------------------------------------------------------
# for TEXT type
DEFAULT_TEXT_ROWS = 4
DEFAULT_TEXT_COLS = 6
# for GOODS type
DEFAULT_GOODS_ROWS = 3
DEFAULT_GOODS_COLS = 4

SLIDE_W, SLIDE_H = 10, 7.5
# rect of one family
FAMILY_W, FAMILY_H = 1.2, 1.5
# Rect of TextBox in a family
TEXT_RECT_SIZE = (0.4, 0.4)
MARGIN_CONNECTOR = 0.01

MARGIN_W = round((SLIDE_W-COLS*FAMILY_W)/(COLS+1), 2)
MARGIN_H = round((SLIDE_H-ROWS*FAMILY_H)/(ROWS+1), 2)
BODY_RECT_TOP = MARGIN_H
BODY_RECT_LEFT = MARGIN_W

IMG_LB_RT = "res/LB-RT.png"
IMG_LT_RB = "res/LT-RB.png"

CONNECTOR_LINE_ADJUST = 0.05

RESULT_RATIO_MAP = list(range(0, 100))
_number_keys = sorted(RESULT_WEIGHT.keys())
NUMBER_MIN = _number_keys[0]
NUMBER_MAX = _number_keys[-1]

output_dir = "output"
output_filename = output_dir + "/add-"+time.strftime("%Y%m%d%H%M%S", time.localtime())+".pptx"
# print(output_filename)

#
# Variables for goods
#
GOODS_FAMILY_W, GOODS_FAMILY_H = 2.6, 0.8
GOODS_MARGIN_W = round((SLIDE_W-COLS*GOODS_FAMILY_W)/(COLS+1), 2)
GOODS_MARGIN_H = round((SLIDE_H-ROWS*GOODS_FAMILY_H)/(ROWS+1), 2)
GOODS_BODY_RECT_TOP = GOODS_MARGIN_H
GOODS_BODY_RECT_LEFT = GOODS_MARGIN_W
GOODS_IMG_PIC_W, GOODS_IMG_PIC_H = 0.25, 0.25
# Layout the pictures in Horizon then Vertical. Ignore the height of container currently
GOODS_IMG_CONTAINER_W, GOODS_IMG_CONTAINER_H = GOODS_FAMILY_W/2.2, 0.0
GOODS_IMG_CONTAINER_PADDING_W, GOODS_IMG_CONTAINER_PADDING_H = 0.0, 0.0
GOODS_IMG_PATH = [
    "res/fruit/apple01.png",
    "res/fruit/apple02.png",
    "res/fruit/banana01.png",
    "res/fruit/banana02.png",
    "res/fruit/cake01.png",
    "res/fruit/cake02.png",
    "res/fruit/cake03.png",
    "res/fruit/cherry01.png",
    "res/fruit/cherry02.png",
    "res/fruit/ice-cream01.png",
    "res/fruit/ice-cream02.png",
    "res/fruit/ice-cream03.png",
    "res/fruit/ice-cream04.png"
]

prs = Presentation()


def slide0():
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Hello, World!"
    subtitle.text = "python-pptx was here!"
    prs.save(output_filename)


def slide1():
    bullet_slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Adding a Bullet Slide'

    tf = body_shape.text_frame
    tf.text = 'Find the bullet slide layout'

    p = tf.add_paragraph()
    p.text = 'Use _TextFrame.text for first bullet'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
    p.level = 2

    prs.save(output_filename)


def slide2():
    img_path = 'res/kongfu.jpg'

    blank_slide_layout = prs.slide_layouts[2]
    slide = prs.slides.add_slide(blank_slide_layout)

    # prs.slide_layouts.remove(blank_slide_layout)
    # slide = prs.slides.add_slide(blank_slide_layout)

    # print(slide.shapes._spTree)
    # for i, j in enumerate(slide.shapes._spTree):
    #     print(i, ', ', j)
    # _shape0 = slide.shapes._spTree[0]
    # _shape1 = slide.shapes._spTree[1]
    # slide.shapes._spTree.remove(_shape0)
    # slide.shapes._spTree.remove(_shape1)

    _clean_default_placeholders(slide)

    left = top = Inches(1)
    height = Inches(1)
    pic = slide.shapes.add_picture(img_path, left, top, height=height)

    left = Inches(5)
    height = Inches(1.5)
    pic = slide.shapes.add_picture(img_path, left, top, height=height)

    prs.save(output_filename)


def new_slides(**kwargs):
    """
    Create PPTX
    """
    for i in range(len(PAGE_CONTENT_STYLE)):
        page_conf = PAGE_CONTENT_STYLE[i]
        op = page_conf['op']
        if op == 'add':
            new_add(i, page_conf, **kwargs)
        elif op == 'minus':
            new_add(i, page_conf, **kwargs)
        elif op == 'multi':
            pass
        elif op == 'division':
            pass
        else:
            raise Exception('Unsupported operation !')


def new_add(slide_index, page_conf, **kwargs):
    blank_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank_slide_layout)
    _clean_default_placeholders(slide)
    _draw_side(slide_index, slide, page_conf, **kwargs)


def _draw_side(slide_index, slide, page_conf, **kwargs):
    max_result = NUMBER_MAX if 'result_upper_limit' not in page_conf else page_conf['result_upper_limit']
    max_result = min(max_result, NUMBER_MAX)
    print('slide {}, lower_limit is {}, result_upper_limit is {}'.format(slide_index, RESULT_LOWER_LIMIT, max_result))
    for i in range(0, ROWS):
        for j in range(0, COLS):
            index_a = RESULT_RATIO_MAP.index(RESULT_LOWER_LIMIT)
            index_b = RESULT_RATIO_MAP.index(max_result)
            result = _rand_result(kwargs['result_min'], index_a, index_b)
            factor1 = _rand_integer((kwargs['factor_min'], result))
            factor2 = result - factor1
            print('[{}][{}], result = {}'.format(i, j, result))
            fn_draw_family = kwargs['fn_draw_family']
            body_rect_top = kwargs['body_rect_top']
            body_rect_left = kwargs['body_rect_left']
            family_h = kwargs['family_h']
            family_w = kwargs['family_w']
            margin_h = kwargs['margin_h']
            margin_w = kwargs['margin_w']
            fn_draw_family(page_conf['result_on_top'], page_conf['op'], factor1, factor2, result,
                           (body_rect_top + (family_h + margin_h) * i, body_rect_left + (family_w + margin_w) * j),
                           slide
                           )
            # _draw_family(
            #     page_conf['result_on_top'], page_conf['op'], factor1, factor2, result,
            #     (BODY_RECT_TOP + (FAMILY_H + MARGIN_H) * i, BODY_RECT_LEFT + (FAMILY_W + MARGIN_W) * j),
            #     slide
            # )
    _draw_textbox(slide, SLIDE_W-0.3, SLIDE_H-0.3, 0.2, 0.2, slide_index+1, RGBColor(128, 128, 128), True)


# --------------------------------------------------------------
# CALLBACK: DRAW TEXT FORMAT SHAPES
# --------------------------------------------------------------
def _cb_draw_family(resultOnTop, op, factor1, factor2, result, pos, slide):
    family_top = pos[0]
    family_left = pos[1]

    if op == 'add' or op == 'minus':
        # Draw result , left and right factor
        active_left = True if random.random() <= 0.5 else False
        if resultOnTop is True:
            tx_width = TEXT_RECT_SIZE[0]
            tx_height = TEXT_RECT_SIZE[1]
            tx_left = family_left + (FAMILY_W - tx_width) / 2
            tx_top = family_top
            left_tx_width = TEXT_RECT_SIZE[0]
            left_tx_height = TEXT_RECT_SIZE[1]
            left_tx_left = family_left
            left_tx_top = family_top + (FAMILY_H - tx_height)
            right_tx_width = TEXT_RECT_SIZE[0]
            right_tx_height = TEXT_RECT_SIZE[1]
            right_tx_left = family_left + (FAMILY_W - tx_width)
            right_tx_top = family_top + (FAMILY_H - tx_height)

            # location for connector line
            line_end_x = family_left + FAMILY_W / 2
            line_end_y = family_top + tx_height + CONNECTOR_LINE_ADJUST
            line_left_begin_x = family_left + tx_width/2
            line_left_begin_y = family_top + (FAMILY_H - tx_height) - CONNECTOR_LINE_ADJUST
            line_right_begin_x = family_left + (FAMILY_W - tx_width/2)
            line_right_begin_y = family_top + (FAMILY_H - tx_height) - CONNECTOR_LINE_ADJUST
        else:
            tx_width = TEXT_RECT_SIZE[0]
            tx_height = TEXT_RECT_SIZE[1]
            tx_left = family_left + (FAMILY_W - tx_width) / 2
            tx_top = family_top + (FAMILY_H - tx_height)
            left_tx_width = TEXT_RECT_SIZE[0]
            left_tx_height = TEXT_RECT_SIZE[1]
            left_tx_left = family_left
            left_tx_top = family_top
            right_tx_width = TEXT_RECT_SIZE[0]
            right_tx_height = TEXT_RECT_SIZE[1]
            right_tx_left = family_left + (FAMILY_W - tx_width)
            right_tx_top = family_top

            # location for connector line
            line_end_x = family_left + FAMILY_W / 2
            line_end_y = family_top + (FAMILY_H - tx_height) - CONNECTOR_LINE_ADJUST
            line_left_begin_x = family_left + tx_width / 2
            line_left_begin_y = family_top + tx_height + CONNECTOR_LINE_ADJUST
            line_right_begin_x = family_left + (FAMILY_W - tx_width / 2)
            line_right_begin_y = family_top + tx_height + CONNECTOR_LINE_ADJUST

        # connector line
        _draw_straight_line(slide, line_left_begin_x, line_left_begin_y,
                            line_end_x - CONNECTOR_LINE_ADJUST, line_end_y)
        _draw_straight_line(slide, line_right_begin_x, line_right_begin_y,
                            line_end_x + CONNECTOR_LINE_ADJUST, line_end_y)

        txt_result, txt_left, txt_right = '', factor1, factor2
        if op == 'minus':
            txt_result = result
            if active_left is True:
                txt_right = ''
            else:
                txt_left = ''

        # result
        _draw_textbox(slide, tx_left, tx_top, tx_width, tx_height, txt_result)
        # left
        _draw_textbox(slide, left_tx_left, left_tx_top, left_tx_width, left_tx_height, txt_left)
        # right
        _draw_textbox(slide, right_tx_left, right_tx_top, right_tx_width, right_tx_height, txt_right)
    elif op == 'multi':
        pass
    elif op == 'division':
        pass


# --------------------------------------------------------------
# CALLBACK: DRAW GOODS FORMAT SHAPES
# --------------------------------------------------------------
def _cb_draw_family_goods(resultOnTop, op, factor1, factor2, result, pos, slide):
    """
    factor1 always has value, and never be 0 forever. We always show factor1 for goods
    """
    family_top = pos[0]
    family_left = pos[1]
    img_path = GOODS_IMG_PATH[random.randint(0, len(GOODS_IMG_PATH)-1)]

    if op == 'add' or op == 'minus':
        # Draw result , left and right factor
        show_factor1_on_left = True if random.random() <= 0.5 else False
        if resultOnTop is True:
            result_anchor_x = family_left + GOODS_FAMILY_W / 2
            result_anchor_y = family_top

            left_child_anchor_x = family_left + GOODS_FAMILY_W / 4
            left_child_anchor_y = family_top + GOODS_FAMILY_H

            right_child_anchor_x = family_left + GOODS_FAMILY_W * 3 / 4
            right_child_anchor_y = family_top + GOODS_FAMILY_H

            # location for connector line
            line_end_x = result_anchor_x
            line_end_y = result_anchor_y + CONNECTOR_LINE_ADJUST
            line_left_begin_x = left_child_anchor_x
            line_left_begin_y = left_child_anchor_y - CONNECTOR_LINE_ADJUST
            line_right_begin_x = right_child_anchor_x
            line_right_begin_y = right_child_anchor_y - CONNECTOR_LINE_ADJUST

        else:
            result_anchor_x = family_left + GOODS_FAMILY_W / 2
            result_anchor_y = family_top + GOODS_FAMILY_H

            left_child_anchor_x = family_left + GOODS_FAMILY_W / 4
            left_child_anchor_y = family_top

            right_child_anchor_x = family_left + GOODS_FAMILY_W * 3 / 4
            right_child_anchor_y = family_top

            # location for connector pic
            line_end_x = result_anchor_x
            line_end_y = result_anchor_y - CONNECTOR_LINE_ADJUST
            line_left_begin_x = left_child_anchor_x
            line_left_begin_y = left_child_anchor_y + CONNECTOR_LINE_ADJUST
            line_right_begin_x = right_child_anchor_x
            line_right_begin_y = right_child_anchor_y + CONNECTOR_LINE_ADJUST

        # connector line
        _draw_straight_line(slide, line_left_begin_x, line_left_begin_y,
                            line_end_x - CONNECTOR_LINE_ADJUST, line_end_y)
        _draw_straight_line(slide, line_right_begin_x, line_right_begin_y,
                            line_end_x + CONNECTOR_LINE_ADJUST, line_end_y)

        txt_result, txt_left, txt_right = '', '', ''
        if op == 'minus':
            txt_result = result
            if show_factor1_on_left is True:
                txt_left = factor1
                txt_right = ''
            else:
                txt_left = ''
                txt_right = factor1
        else:
            if show_factor1_on_left is True:
                txt_left = factor1
                txt_right = factor2
            else:
                txt_left = factor2
                txt_right = factor1

        # result
        _draw_goods(slide, number=txt_result,
                    anchor_x=result_anchor_x, anchor_y=result_anchor_y,
                    anchor_side="bottom_middle" if resultOnTop is True else "top_middle",
                    is_text=False if txt_result == result else True,
                    img=img_path)
        # left
        _draw_goods(slide, number=txt_left,
                    anchor_x=left_child_anchor_x, anchor_y=left_child_anchor_y,
                    anchor_side="top_middle" if resultOnTop is True else "bottom_middle",
                    is_text=True if (txt_left == '' or txt_left == 0) else False,
                    img=img_path)
        # right
        _draw_goods(slide, number=txt_right,
                    anchor_x=right_child_anchor_x, anchor_y=right_child_anchor_y,
                    anchor_side="top_middle" if resultOnTop is True else "bottom_middle",
                    is_text=True if (txt_right == '' or txt_right == 0) else False,
                    img=img_path)
    elif op == 'multi':
        pass
    elif op == 'division':
        pass


def _draw_goods(slide, number, anchor_x, anchor_y, anchor_side, is_text, img):
    if is_text is True:
        tx_width = TEXT_RECT_SIZE[0]
        tx_height = TEXT_RECT_SIZE[1]
        if anchor_side == "top_middle":
            tx_left = anchor_x - tx_width/2
            tx_top = anchor_y
        elif anchor_side == "bottom_middle":
            tx_left = anchor_x - tx_width/2
            tx_top = anchor_y - tx_height
        _draw_textbox(slide, tx_left, tx_top, tx_width, tx_height, number)
    else:
        if anchor_side == "top_middle":
            # anchor stands at top-middle
            #         (x,y)
            # ----------o---------
            # |                  |
            # |                  |
            # |                  |
            #          ...
            #
            fromPos = {'x': anchor_x - GOODS_IMG_CONTAINER_W/2, 'y': anchor_y}
            _draw_goods_img(slide, number, fromPos, 'down',
                            GOODS_IMG_CONTAINER_PADDING_W, GOODS_IMG_CONTAINER_PADDING_H,
                            img, GOODS_IMG_PIC_W, GOODS_IMG_PIC_H)
        elif anchor_side == "bottom_middle":
            # anchor stands at bottom-middle
            #          ...
            # |                  |
            # |                  |
            # |                  |
            # ----------o---------
            #         (x,y)
            #
            fromPos = {'x': anchor_x - GOODS_IMG_CONTAINER_W / 2, 'y': anchor_y}
            _draw_goods_img(slide, number, fromPos, 'up',
                            GOODS_IMG_CONTAINER_PADDING_W, GOODS_IMG_CONTAINER_PADDING_H,
                            img, GOODS_IMG_PIC_W, GOODS_IMG_PIC_H)


def _draw_goods_img(slide, number, fromPos, layoutDirect, padding_w, padding_h, img, img_w, img_h):
    per_row = int(GOODS_IMG_CONTAINER_W / img_w)
    if number < per_row:
        new_left = fromPos['x'] + (GOODS_IMG_CONTAINER_W - number*img_w - (number-1)*padding_w)/2
        new_top = fromPos['y']
        i = 0
        for i in range(0, number):
            pic_left = new_left + (img_w + padding_w) * i
            if layoutDirect == "down":
                pic_top = new_top
            else:
                pic_top = new_top - img_h
            _draw_pic(slide, Inches(pic_left), Inches(pic_top), Inches(img_h), img)
    else:
        row, col = math.ceil(number/per_row), per_row
        for i in range(0, row):
            for j in range(0, col):
                if (i*col+j) >= number:
                    break
                pic_left = fromPos['x'] + (img_w + padding_w)*j
                if layoutDirect == "down":
                    pic_top = fromPos['y'] + (img_h + padding_h)*i
                else:
                    pic_top = fromPos['y'] - img_h - (img_h + padding_h) * i
                _draw_pic(slide, Inches(pic_left), Inches(pic_top), Inches(img_h), img)


def test():
    blank_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank_slide_layout)
    _clean_default_placeholders(slide)
    # _draw_goods(slide, number=3, anchor_x=0, anchor_y=0, anchor_side="top_middle", is_text=False, img=GOODS_IMG_PATH[0])


# --------------------------------------------------------------
# Common functions
# --------------------------------------------------------------
def _draw_textbox(slide, left, top, width, height, text, text_color=RGBColor(0, 0, 0), dash=False):
    txBox = slide.shapes.add_textbox(
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height))
    txBox.text = str(text)
    if dash is True:
        txBox.line.dash_style = MSO_LINE.ROUND_DOT
    txBox.line.color.rgb = RGBColor(0, 0, 0)
    txBox.line.width = Pt(1.0)
    txBox.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    txBox.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    txBox.text_frame.paragraphs[0].font.color.rgb = text_color
    # txBox.fill.solid()
    # txBox.fill.fore_color.rgb = RGBColor(0, 0, 0)


def _draw_pic(slide, left, top, height, img):
    pic = slide.shapes.add_picture(img, left, top, height=height)


def _draw_connector(slide, loc_from, loc_to, height, img, alignLoc):
    _rect = [loc_to['x'] - loc_from['x'], abs(loc_to['y'] - loc_from['y'])]

    pic_left, pic_top = Inches(0), Inches(0)
    pic_height = Inches(height)

    if loc_to['y'] < loc_from['y']:
        if alignLoc == "from":
            pic_left, pic_top = Inches(loc_from['x']), Inches(loc_from['y'] - _rect[1])
        elif alignLoc == "to":
            pic_left, pic_top = Inches(loc_to['x'] - _rect[0] - 0.1), Inches(loc_to['y'])
    else:
        if alignLoc == "from":
            pic_left, pic_top = Inches(loc_from['x']), Inches(loc_from['y'])
        elif alignLoc == "to":
            pic_left, pic_top = Inches(loc_to['x'] - _rect[0] - 0.1), Inches(loc_to['y'] - _rect[1])
    # pic = slide.shapes.add_picture(img, pic_left, pic_top, height=pic_height)
    _draw_pic(slide, pic_left, pic_top, pic_height, img)


def _draw_straight_line(slide, begin_x, begin_y, end_x, end_y):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(begin_x), Inches(begin_y), Inches(end_x), Inches(end_y)
    )
    connector.shadow.inherit = False
    connector.line.color.rgb = RGBColor(0, 0, 0)
    connector.line.width = Pt(1.0)


def _rand_result(lower_limit, index_a, index_b):
    _index = random.randint(index_a, index_b)
    rslt = RESULT_RATIO_MAP[_index]
    return rslt if rslt >= lower_limit else _rand_result(lower_limit, index_a, index_b)


def _rand_integer(rng):
    return random.randint(rng[0], rng[1])


def _clean_default_placeholders(slide):
    textbox = slide.shapes[0]
    sp = textbox.element
    sp.getparent().remove(sp)
    textbox = slide.shapes[0]
    sp = textbox.element
    sp.getparent().remove(sp)


def _init():
    global RESULT_WEIGHT
    totalweight = 0.0
    number_keys = sorted(RESULT_WEIGHT.keys())
    min_number = number_keys[0]
    max_number = number_keys[-1]
    prev_weight = RESULT_WEIGHT[number_keys[0]]
    while min_number <= max_number:
        if min_number in number_keys:
            prev_weight = RESULT_WEIGHT[min_number]
        else:
            RESULT_WEIGHT[min_number] = prev_weight
        totalweight += prev_weight
        min_number += 1

    # sort dictionary RESULT_WEIGHT
    _sorted_result_weight = {}
    ss = sorted(RESULT_WEIGHT.items())
    [_sorted_result_weight.update({k: v}) for k, v in ss]
    RESULT_WEIGHT = _sorted_result_weight

    acc_ratio = 0
    cur_index = 0
    for k, v in RESULT_WEIGHT.items():
        ratio_in_100 = max(1, int((v * 100)/totalweight))
        acc_ratio += ratio_in_100
        print("number = {}, ratio_in_100 = {}".format(k, ratio_in_100))
        # if k == 30:
        #     print('')
        if acc_ratio > 100:
            ratio_in_100 -= (acc_ratio - 100)
        if k == len(RESULT_WEIGHT)-1 and acc_ratio < 100:
            ratio_in_100 += 100 - acc_ratio
        i = 0
        while i < ratio_in_100:
            # print("{} -> {}".format(cur_index, k))
            RESULT_RATIO_MAP[cur_index] = k
            cur_index += 1
            i += 1


def usage():
    print('- Program: math.py')
    print('- Author:  liuhao')
    print('- Copyright © 2020')
    print('')
    print('math.py -t <test type>')
    print('command:')
    print('-t --type      text : display text for number')
    print('               goods : display fruit for number')


if __name__ == "__main__":
    # test()
    number_type = ''
    try:
        opts, args = getopt.getopt(sys.argv[1:], "ht:", ["help", "type="])
    except getopt.GetoptError:
        usage()
        sys.exit(2)
    for opt, arg in opts:
        if opt in ("-h", "--help"):
            usage()
            sys.exit()
        elif opt in ("-t", "--type"):
            number_type = arg
        else:
            usage()
            sys.exit()

    if len(opts) == 0:
        usage()
        sys.exit()
    _init()
    if number_type == 'text':
        new_slides(
            fn_draw_family=_cb_draw_family,
            body_rect_top=BODY_RECT_TOP, body_rect_left=BODY_RECT_LEFT,
            family_h=FAMILY_H, family_w=FAMILY_W,
            margin_h=MARGIN_H, margin_w=MARGIN_W,
            result_min=1, factor_min=0
        )
    elif number_type == 'goods':
        new_slides(
            fn_draw_family=_cb_draw_family_goods,
            body_rect_top=GOODS_BODY_RECT_TOP, body_rect_left=GOODS_BODY_RECT_LEFT,
            family_h=GOODS_FAMILY_H, family_w=GOODS_FAMILY_W,
            margin_h=GOODS_MARGIN_H, margin_w=GOODS_MARGIN_W,
            result_min=1, factor_min=1
        )
    prs.save(output_filename)
