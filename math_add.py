import time
import random
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT as PP_ALIGN
from pptx.enum.text import MSO_VERTICAL_ANCHOR as MSO_ANCHOR
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MSO_SHAPE

'''
1. Remove placeholders in slide
https://python-pptx.readthedocs.io/en/latest/user/placeholders-understanding.html
https://stackoverflow.com/questions/39603318/how-to-delete-unpopulated-placeholder-items-using-python-pptx

2. Specifies the horizontal alignment for one or more paragraphs.
https://python-pptx.readthedocs.io/en/latest/api/text.html#paragraph-objects
https://python-pptx.readthedocs.io/en/latest/api/enum/PpParagraphAlignment.html#ppparagraphalignment

'''

"""
Side Rect (Inches):
       10
+--------------+
|              |    7.5 
|              |
+--------------+
Width & Height of one slide is about 10 & 7.5
"""
# const values
SLIDE_W, SLIDE_H = 10, 7.5
# rows and columns
ROWS = 3
COLS = 6
# rect of one family
FAMILY_W, FAMILY_H = 0.6, 0.8
# Rect of TextBox in a family
TEXT_RECT_SIZE = (0.4, 0.4)

MARGIN_W = round((SLIDE_W-COLS*FAMILY_W)/(COLS+1), 2)
MARGIN_H = round((SLIDE_H-ROWS*FAMILY_H)/(ROWS+1), 2)
BODY_RECT_TOP = MARGIN_H
BODY_RECT_LEFT = MARGIN_W

# ratio of random number
RESULT_WEIGHT = {
    0: 1,
    1: 1,
    2: 1,
    3: 1,
    4: 1,
    5: 1,
    6: 3,
    7: 3,
    8: 3,
    9: 3,
    10: 3
}

RESULT_RATIO_MAP = list(range(0, 100))

prs = Presentation()

output_dir = "output"
output_filename = output_dir + "/add-"+time.strftime("%Y%m%d%H%M%S", time.localtime())+".pptx"
print(output_filename)


def slide0():
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Hello, World!"
    subtitle.text = "python-pptx was here!"
    prs.save(output_filename)


def slide1():
    bullet_slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Adding a Bullet Slide'

    tf = body_shape.text_frame
    tf.text = 'Find the bullet slide layout'

    p = tf.add_paragraph()
    p.text = 'Use _TextFrame.text for first bullet'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
    p.level = 2

    prs.save(output_filename)


def slide2():
    img_path = 'res/kongfu.jpg'

    blank_slide_layout = prs.slide_layouts[2]
    slide = prs.slides.add_slide(blank_slide_layout)

    # prs.slide_layouts.remove(blank_slide_layout)
    # slide = prs.slides.add_slide(blank_slide_layout)

    # print(slide.shapes._spTree)
    # for i, j in enumerate(slide.shapes._spTree):
    #     print(i, ', ', j)
    # _shape0 = slide.shapes._spTree[0]
    # _shape1 = slide.shapes._spTree[1]
    # slide.shapes._spTree.remove(_shape0)
    # slide.shapes._spTree.remove(_shape1)

    _clean_default_placeholders(slide)

    left = top = Inches(1)
    height = Inches(1)
    pic = slide.shapes.add_picture(img_path, left, top, height=height)

    left = Inches(5)
    height = Inches(1.5)
    pic = slide.shapes.add_picture(img_path, left, top, height=height)

    prs.save(output_filename)


def new_slides(nm, op, upper_limit):
    """
    Create PPTX

    :param nm: num of slide
    :param op: add/minus/multi/division
    :param upper_limit: max result of the operation
    :return:
    """

    if op == 'add':
        new_add(nm, upper_limit)
    elif op == 'minus':
        pass
    elif op == 'multi':
        pass
    elif op == 'division':
        pass
    else:
        raise Exception('Unsupported operation !')


def new_add(nm, upper_limit):
    i = 0
    while i < nm:
        blank_slide_layout = prs.slide_layouts[i]
        slide = prs.slides.add_slide(blank_slide_layout)
        _clean_default_placeholders(slide)
        _draw_side(i, slide, upper_limit)
        i += 1


def _draw_side(page, slide, upper_limit):
    print('slide ' + str(page))
    for i in range(0, ROWS):
        # resultOnTop = False if random.randint(0, 1) == 0 else True
        resultOnTop = False
        for j in range(0, COLS):
            result = _rand_result(1)
            factor1 = _rand_integer((0, result))
            factor2 = result - factor1
            print('[{}][{}], result = {}'.format(i, j, result))
            _draw_tb(
                resultOnTop, True, factor1, factor2, result,
                (BODY_RECT_TOP + (FAMILY_H + MARGIN_H) * i, BODY_RECT_LEFT + (FAMILY_W + MARGIN_W) * j),
                slide
            )


def _draw_tb(resultOnTop, isAdd, factor1, factor2, result, pos, slide):
    family_top = pos[0]
    family_left = pos[1]

    if isAdd is True:
        if resultOnTop is True:
            # add textbox
            tx_width = TEXT_RECT_SIZE[0]
            tx_height = TEXT_RECT_SIZE[1]
            tx_left = family_left + (FAMILY_W - tx_width)/2
            tx_top = family_top
        else:
            # add textbox
            tx_width = TEXT_RECT_SIZE[0]
            tx_height = TEXT_RECT_SIZE[1]
            tx_left = family_left + (FAMILY_W - tx_width) / 2
            tx_top = family_top + (FAMILY_H - tx_height)
        # print("text box rect: {}, {}, {}, {}".format(tx_left, tx_top, tx_width, tx_height))
        txBox = slide.shapes.add_textbox(
            Inches(tx_left),
            Inches(tx_top),
            Inches(tx_width),
            Inches(tx_height))
        txBox.text = str(result)
        txBox.line.color.rgb = RGBColor(0, 0, 0)
        txBox.line.width = Pt(1.0)
        txBox.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        txBox.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        # txBox.fill.solid()
        # txBox.fill.fore_color.rgb = RGBColor(0, 0, 0)
    else:
        pass


def _rand_result(minnum):
    rand_index = random.randint(0, 99)
    rslt = RESULT_RATIO_MAP[rand_index]
    # if rslt < minnum:
    #     rslt = _rand_result(minnum)
    return rslt if rslt >= minnum else _rand_result(minnum)


def _rand_integer(rng):
    return random.randint(rng[0], rng[1])


def _clean_default_placeholders(slide):
    textbox = slide.shapes[0]
    sp = textbox.element
    sp.getparent().remove(sp)
    textbox = slide.shapes[0]
    sp = textbox.element
    sp.getparent().remove(sp)


def _init():
    totalweight = 0
    for w in RESULT_WEIGHT.values():
        totalweight += w
    acc_ratio = 0
    cur_index = 0
    for k, v in RESULT_WEIGHT.items():
        # print("{} ==== {}".format(k, v))
        ratio_in_100 = round((v * 100)/totalweight)
        # print("ratio_in_100 = {}".format(ratio_in_100))
        acc_ratio += ratio_in_100
        if acc_ratio > 100:
            ratio_in_100 -= (acc_ratio - ratio_in_100)
        i = 0
        while i < ratio_in_100:
            # print("{} -> {}".format(cur_index, k))
            RESULT_RATIO_MAP[cur_index] = k
            cur_index += 1
            i += 1


if __name__ == "__main__":
    # slide0()
    # slide1()
    # slide2()

    _init()
    new_slides(2, "add", 10)
    prs.save(output_filename)
